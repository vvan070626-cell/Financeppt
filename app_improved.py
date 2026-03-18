import streamlit as st
from langchain_openai import ChatOpenAI
from langchain_core.messages import SystemMessage, HumanMessage, AIMessage
from pptx import Presentation
import fitz  # PyMuPDF
import io
import sqlite3
from datetime import datetime

# ─────────────────────────────────────────────
# 页面配置
# ─────────────────────────────────────────────
st.set_page_config(
    page_title="AI 学习助手",
    page_icon="[AI]",
    layout="wide"
)

# ─────────────────────────────────────────────
# 初始化 LLM
# ─────────────────────────────────────────────
api_key = st.secrets["DEEPSEEK_API_KEY"]

llm = ChatOpenAI(
    model="deepseek-chat",
    openai_api_key=api_key,
    base_url="https://api.deepseek.com",
    max_retries=2
)

# ─────────────────────────────────────────────
# 文本提取函数
# ─────────────────────────────────────────────
def extract_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join([run.text for run in para.runs]).strip()
                    if line:
                        texts.append(line)
    return "\n".join(texts)

def extract_pdf(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    return "\n".join([page.get_text() for page in doc])

def extract_content(uploaded_file):
    file_bytes = uploaded_file.read()
    name = uploaded_file.name.lower()
    if name.endswith(".pptx"):
        return extract_pptx(file_bytes)
    elif name.endswith(".pdf"):
        return extract_pdf(file_bytes)
    return None

# ─────────────────────────────────────────────
# SQLite 工具函数
# ─────────────────────────────────────────────
DB_PATH = "ai_study_assistant.sqlite3"

def db_conn():
    return sqlite3.connect(DB_PATH, check_same_thread=False)

def db_init():
    con = db_conn()
    cur = con.cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS files (
            filename   TEXT PRIMARY KEY,
            content    TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )
    """
    )
    cur.execute("""
        CREATE TABLE IF NOT EXISTS messages (
            id         INTEGER PRIMARY KEY AUTOINCREMENT,
            filename   TEXT NOT NULL,
            role       TEXT NOT NULL,
            content    TEXT NOT NULL,
            created_at TEXT NOT NULL
        )
    """
    )
    con.commit()
    con.close()

def db_upsert_file(filename: str, content: str):
    con = db_conn()
    cur = con.cursor()
    cur.execute("""
        INSERT INTO files(filename, content, updated_at)
        VALUES (?, ?, ?)
        ON CONFLICT(filename) DO UPDATE SET
            content=excluded.content,
            updated_at=excluded.updated_at
    """, (filename, content, datetime.utcnow().isoformat()))
    con.commit()
    con.close()

def db_add_message(filename: str, role: str, content: str):
    con = db_conn()
    cur = con.cursor()
    cur.execute("""
        INSERT INTO messages(filename, role, content, created_at)
        VALUES (?, ?, ?, ?)
    """, (filename, role, content, datetime.utcnow().isoformat()))
    con.commit()
    con.close()

def db_load_all():
    """从数据库恢复所有文件内容和聊天记录"""
    con = db_conn()
    cur = con.cursor()

    cur.execute("SELECT filename, content FROM files ORDER BY updated_at DESC")
    file_contents = {fn: ct for fn, ct in cur.fetchall()}

    chat_histories = {fn: [] for fn in file_contents.keys()}
    cur.execute("SELECT filename, role, content FROM messages ORDER BY id ASC")
    for fn, role, content in cur.fetchall():
        if fn in chat_histories:
            chat_histories[fn].append({"role": role, "content": content})

    con.close()
    return file_contents, chat_histories

def db_clear_chat(filename: str):
    con = db_conn()
    cur = con.cursor()
    cur.execute("DELETE FROM messages WHERE filename=?", (filename,))
    con.commit()
    con.close()

def db_delete_file(filename: str):
    """从数据库删除文件及其对话记录"""
    con = db_conn()
    cur = con.cursor()
    cur.execute("DELETE FROM messages WHERE filename=?", (filename,))
    cur.execute("DELETE FROM files WHERE filename=?", (filename,))
    con.commit()
    con.close()

# ─────────────────────────────────────────────
# 启动时：初始化 DB，并从 DB 恢复状态
# ─────────────────────────────────────────────
db_init()

if "file_contents" not in st.session_state:
    st.session_state.file_contents = {}

if "chat_histories" not in st.session_state:
    st.session_state.chat_histories = {}

# 关键：刷新/重启后，如果 session_state 是空的，就从 DB 恢复
if not st.session_state.file_contents and not st.session_state.chat_histories:
    fc, ch = db_load_all()
    st.session_state.file_contents = fc
    st.session_state.chat_histories = ch

# ─────────────────────────────────────────────
# 侧边栏：文件上传
# ─────────────────────────────────────────────
with st.sidebar:
    st.title("AI 学习助手")
    st.caption("上传 PPT / PDF，在 Tabs 中独立对话")

    st.divider()

    # 文件上传区
    uploaded_files = st.file_uploader(
        "上传文件（支持 .pptx / .pdf）",
        type=["pptx", "pdf"],
        accept_multiple_files=True
    )

    # 处理新上传的文件
    if uploaded_files:
        for file in uploaded_files:
            fname = file.name
            if fname not in st.session_state.file_contents:
                with st.spinner(f"正在提取 {fname}..."):
                    content = extract_content(file)
                    if content:
                        st.session_state.file_contents[fname] = content
                        st.session_state.chat_histories[fname] = []
                        db_upsert_file(fname, content)   # 持久化到 DB
                        st.success(f"{fname} 提取成功！")
                    else:
                        st.error(f"{fname} 提取失败，请检查格式。")

# ─────────────────────────────────────────────
# 主区域：多 Tab 聊天窗口
# ─────────────────────────────────────────────
if st.session_state.file_contents:
    # 获取所有文件名
    filenames = list(st.session_state.file_contents.keys())
    
    # 创建 tabs，每个文件一个 tab
    tabs = st.tabs(filenames)
    
    # 为每个 tab 创建对话界面
    for tab, fname in zip(tabs, filenames):
        with tab:
            content = st.session_state.file_contents[fname]
            history = st.session_state.chat_histories[fname]
            
            # 标题栏 + 清空/删除按钮
            col_title, col_buttons = st.columns([7, 2])
            with col_title:
                st.markdown(f"### 📄 {fname}")
            with col_buttons:
                col_clear, col_delete = st.columns(2)
                with col_clear:
                    if history:
                        if st.button("🗑️ 清空", key=f"clear_{fname}"):
                            st.session_state.chat_histories[fname] = []
                            db_clear_chat(fname)
                            st.rerun()
                with col_delete:
                    if st.button("✕ 删除", key=f"del_{fname}"):
                        db_delete_file(fname)
                        del st.session_state.file_contents[fname]
                        del st.session_state.chat_histories[fname]
                        st.rerun()
            
            st.divider()
            
            # 显示历史对话
            # 使用 container 来实现可滚动的聊天区域
            chat_container = st.container()
            with chat_container:
                for msg in history:
                    with st.chat_message(msg["role"]):
                        st.write(msg["content"])
            
            st.divider()
            
            # 用户输入框（放在底部）
            user_input = st.chat_input(
                placeholder=f"问关于 {fname} 的问题...",
                key=f"chat_input_{fname}"
            )

            if user_input:
                # 追加用户消息到 session_state 和 DB
                history.append({"role": "user", "content": user_input})
                db_add_message(fname, "user", user_input)

                with st.chat_message("user"):
                    st.write(user_input)

                # 构建发送给 LLM 的消息列表
                messages = [
                    SystemMessage(content=(
                        f"你是一位专业的学习助手，擅长分析学科知识点。
                        f"以下是用户上传的文档内容，请基于这份文档回答用户的所有问题：
                        f"\n
                        f"【文档内容】\n{content}\n
                        f"请用中文回答，回答要清晰、结构化。"
                    ))
                ]

                # 多轮对话记忆
                for msg in history:
                    if msg["role"] == "user":
                        messages.append(HumanMessage(content=msg["content"]))
                    else:
                        messages.append(AIMessage(content=msg["content"]))

                # 调用 AI
                with st.chat_message("assistant"):
                    with st.spinner("AI 思考中..."):
                        try:
                            response = llm.invoke(messages)
                            reply = response.content
                            st.write(reply)

                            # 追加 AI 回复到 session_state 和 DB
                            history.append({"role": "assistant", "content": reply})
                            db_add_message(fname, "assistant", reply)

                            # 写回 session_state
                            st.session_state.chat_histories[fname] = history

                        except Exception as e:
                            st.error(f"API 调用失败：{str(e)}")

else:
    # 没有任何文件时，显示引导提示
    st.title("AI 学习助手")
    st.info("请从左侧栏上传至少一个 .pptx 或 .pdf 文件以开始对话。")